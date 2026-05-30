#!/usr/bin/env python3
"""
本地文件浏览器 HTTP 服务
=========================
在浏览器中浏览 workspace 目录结构，支持预览和下载。
预览格式：md → HTML渲染 | 图片 直接显示 | PDF 嵌入预览 |
          docx/pptx/xlsx → HTML预览 | 纯文本直接显示

启动: python3 file-browser.py [--port 8765]
"""

import http.server
import os
import json
import sys
import mimetypes
import urllib.parse
import io
import time
from pathlib import Path

try:
    import markdown as _md_lib
    HAS_MD = True
except ImportError:
    HAS_MD = False

try:
    from docx import Document
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    from pptx import Presentation
    from pptx.util import Inches
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    import openpyxl
    HAS_XLSX = True
except ImportError:
    HAS_XLSX = False

try:
    import pdfplumber
    HAS_PDF = True
except ImportError:
    HAS_PDF = False

# ---------- config ----------
PORT = int(sys.argv[2]) if len(sys.argv) > 2 and sys.argv[1] == '--port' else 8765
ROOT = Path(os.environ.get("FILEBROWSER_ROOT", os.path.join(os.path.dirname(os.path.abspath(__file__)), ""))).resolve()

# 忽略的目录/文件
IGNORE_DIRS = {'.git', '__pycache__', '.DS_Store', 'node_modules', '.venv', 'venv', '.backups', '_backups'}
IGNORE_FILES = {'.DS_Store', 'Thumbs.db'}
# 允许显示的二进制扩展名
IMAGE_EXTS = {'.png', '.jpg', '.jpeg', '.gif', '.webp', '.svg', '.ico', '.bmp'}
OFFICE_EXTS = {'.docx', '.pptx', '.xlsx', '.pdf'}


def is_safe_path(path_str):
    """防止路径穿越"""
    clean = path_str.lstrip('/') if path_str else ''
    try:
        resolved = (ROOT / clean).resolve()
        return resolved.is_relative_to(ROOT)
    except (ValueError, OSError):
        return False


def get_file_info(abs_path, relative_to_root):
    """获取单个文件信息"""
    stat = abs_path.stat()
    rel = str(relative_to_root)
    return {
        'name': abs_path.name,
        'path': '/' + rel,
        'size': stat.st_size,
        'size_human': _human_size(stat.st_size),
        'modified': time.strftime('%Y-%m-%d %H:%M', time.localtime(stat.st_mtime)),
        'is_dir': abs_path.is_dir(),
        'ext': abs_path.suffix.lower(),
    }


def _human_size(size):
    for unit in ['B', 'KB', 'MB', 'GB']:
        if size < 1024:
            return f"{size:.1f} {unit}"
        size /= 1024
    return f"{size:.1f} TB"


def build_tree(rel_dir=''):
    """构建目录树 JSON"""
    target = ROOT / rel_dir.lstrip('/') if rel_dir else ROOT
    target = target.resolve()
    if not target.is_relative_to(ROOT):
        target = ROOT

    items = []
    try:
        for entry in sorted(target.iterdir(), key=lambda e: (not e.is_dir(), e.name.lower())):
            if entry.name in IGNORE_FILES:
                continue
            if entry.is_dir() and entry.name in IGNORE_DIRS:
                continue
            # 跳过隐藏文件/夹
            if entry.name.startswith('.') and entry.name not in {'.gitignore', '.env.example'}:
                continue

            rel = entry.relative_to(ROOT)
            info = get_file_info(entry, rel)
            if entry.is_dir():
                info['children'] = []  # 懒加载
            items.append(info)
    except PermissionError:
        pass
    return items


def build_breadcrumb(path_parts):
    """面包屑HTML"""
    parts = []
    cumulative = ''
    parts.append('<a href="/" class="crumb">🏠</a>')
    for part in path_parts:
        if not part:
            continue
        cumulative += '/' + part
        parts.append(f'<span class="crumb-sep">/</span><a href="/?dir={urllib.parse.quote(cumulative)}" class="crumb">{part}</a>')
    return ''.join(parts)


def render_md_to_html(md_text):
    """Markdown → HTML"""
    if HAS_MD:
        return _md_lib.markdown(md_text, extensions=['fenced_code', 'codehilite', 'tables', 'toc'])
    # 回退：简单处理
    escaped = md_text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
    return f'<pre style="white-space:pre-wrap;font-family:monospace;font-size:0.85em">{escaped}</pre>'


def render_docx_to_html(path):
    """DOCX → HTML"""
    doc = Document(path)
    parts = []
    for para in doc.paragraphs:
        if para.style.name.startswith('Heading'):
            level = para.style.name.split()[-1]
            parts.append(f'<h{level} style="margin:10px 0 4px;color:var(--ink)">{para.text}</h{level}>')
        elif para.text.strip():
            parts.append(f'<p style="margin:4px 0;line-height:1.7">{para.text}</p>')
    # 表格
    for table in doc.tables:
        parts.append('<table class="preview-table">')
        for row in table.rows:
            parts.append('<tr>')
            for cell in row.cells:
                parts.append(f'<td style="padding:4px 8px;border:1px solid var(--border)">{cell.text}</td>')
            parts.append('</tr>')
        parts.append('</table>')
    return '\n'.join(parts) if parts else '<p style="color:var(--ink-light)">(空文档)</p>'


def render_pptx_to_html(path):
    """PPTX → HTML (文本提取)"""
    prs = Presentation(path)
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f'<div style="border:2px solid var(--accent-gold);border-radius:8px;padding:12px;margin:12px 0;background:var(--paper-white)"><h3 style="margin:0 0 8px;color:var(--accent-gold)">📊 幻灯片 {i}</h3>')
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        parts.append(f'<p style="margin:2px 0;line-height:1.6">{text}</p>')
        parts.append('</div>')
    return '\n'.join(parts) if parts else '<p style="color:var(--ink-light)">(无文本内容或空演示文稿)</p>'


def render_xlsx_to_html(path):
    """XLSX → HTML 表格"""
    wb = openpyxl.load_workbook(path, data_only=True)
    parts = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        parts.append(f'<h3 style="margin:16px 0 8px;color:var(--accent-gold)">📋 {sheet_name}</h3>')
        parts.append('<div style="overflow-x:auto"><table class="preview-table">')
        for row in ws.iter_rows(max_row=min(ws.max_row, 500), values_only=True):
            parts.append('<tr>')
            for cell in row:
                val = str(cell) if cell is not None else ''
                parts.append(f'<td style="padding:3px 6px;border:1px solid var(--border);font-size:0.8em;max-width:300px;overflow:hidden;white-space:nowrap;text-overflow:ellipsis">{val}</td>')
            parts.append('</tr>')
        parts.append('</table></div>')
    wb.close()
    return '\n'.join(parts) if parts else '<p style="color:var(--ink-light)">(空工作簿)</p>'


def render_pdf_preview_html(path):
    """PDF → 文本提取"""
    parts = []
    with pdfplumber.open(path) as pdf:
        for i, page in enumerate(pdf.pages, 1):
            text = page.extract_text()
            if text:
                parts.append(f'<div style="border:1px solid var(--border);border-radius:4px;padding:8px;margin:8px 0"><strong style="color:var(--accent-gold)">第 {i} 页</strong><pre style="white-space:pre-wrap;font-size:0.82em;margin:4px 0">{text}</pre></div>')
    return '\n'.join(parts) if parts else '<p style="color:var(--ink-light)">(无法提取文本内容)</p>'


# ---------- HTML 模板 ----------
HTML_TEMPLATE = '''<!DOCTYPE html>
<html lang="zh-CN">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>满意红 · 文件浏览器</title>
<style>
  :root {
    --bg: #F5F0E6;
    --paper-white: #FDFAF3;
    --card: #FFFDF8;
    --ink: #4A3728;
    --ink-light: #8B7355;
    --ink-lighter: #B8A898;
    --accent-red: #C23B22;
    --accent-gold: #B8860B;
    --border: #E0D5C0;
    --border-light: #EDE5D5;
    --hover: #F0E8D8;
    --code-bg: #2D2A26;
    --code-fg: #E8DCC8;
  }
  * { box-sizing:border-box; margin:0; padding:0; }
  body {
    font-family: -apple-system,BlinkMacSystemFont,"Segoe UI",Roboto,"Noto Sans SC",sans-serif;
    background: var(--bg);
    color: var(--ink);
    min-height:100vh;
    display:flex;
  }
  a { color: var(--accent-gold); text-decoration:none; }
  a:hover { text-decoration:underline; }

  /* 侧边栏 */
  #sidebar {
    width: 320px;
    min-width: 260px;
    max-width: 420px;
    background: var(--card);
    border-right: 1px solid var(--border);
    height: 100vh;
    overflow-y: auto;
    resize: horizontal;
    display:flex;
    flex-direction:column;
  }
  #sidebar-header {
    padding: 16px;
    border-bottom: 1px solid var(--border);
    background: linear-gradient(135deg, var(--accent-red), #A01E14);
    color: #FFF;
  }
  #sidebar-header h2 { font-size:1.15em; margin-bottom:4px; }
  #sidebar-header p { font-size:0.72em; opacity:0.85; }
  #sidebar-resizer {
    width:4px;
    cursor:col-resize;
    background:var(--border-light);
    flex-shrink:0;
  }
  #tree-container {
    flex:1;
    overflow-y:auto;
    padding:8px 0;
  }
  .tree-item {
    display:flex;
    align-items:center;
    padding:5px 12px 5px 16px;
    cursor:pointer;
    font-size:0.88em;
    border-left: 3px solid transparent;
    transition: background 0.15s;
    user-select:none;
    white-space:nowrap;
  }
  .tree-item:hover { background:var(--hover); }
  .tree-item.active { border-left-color:var(--accent-gold); background:var(--hover); }
  .tree-item .icon { margin-right:6px; width:18px; text-align:center; flex-shrink:0; }
  .tree-item .name { overflow:hidden; text-overflow:ellipsis; }
  .tree-item .meta { margin-left:auto; font-size:0.7em; color:var(--ink-lighter); flex-shrink:0; }
  .tree-children { padding-left:0; }
  .tree-children .tree-item { padding-left:28px; }
  .tree-children .tree-children .tree-item { padding-left:40px; }
  .tree-children .tree-children .tree-children .tree-item { padding-left:52px; }
  .tree-toggle {
    display:inline-block;
    width:14px;
    font-size:0.65em;
    color:var(--ink-lighter);
    transition:transform 0.2s;
    flex-shrink:0;
    text-align:center;
  }
  .tree-toggle.expanded { transform:rotate(90deg); }
  .tree-toggle.leaf { visibility:hidden; }
  .tree-children.collapsed { display:none; }

  /* 主区域 */
  #main {
    flex:1;
    display:flex;
    flex-direction:column;
    height:100vh;
    overflow:hidden;
  }
  #toolbar {
    padding:10px 20px;
    background:var(--card);
    border-bottom:1px solid var(--border);
    display:flex;
    align-items:center;
    gap:12px;
    flex-shrink:0;
  }
  #breadcrumb { font-size:0.88em; color:var(--ink-light); display:flex; align-items:center; gap:4px; flex-wrap:wrap; }
  .crumb { color:var(--ink-light); }
  .crumb:hover { color:var(--accent-gold); text-decoration:none; }
  .crumb-sep { color:var(--ink-lighter); margin:0 2px; }
  .toolbtn {
    background:var(--accent-gold);
    color:#FFF;
    border:none;
    padding:6px 14px;
    border-radius:6px;
    font-size:0.82em;
    cursor:pointer;
    font-weight:600;
    white-space:nowrap;
  }
  .toolbtn:hover { background:#9A6F09; }
  .toolbtn.plain {
    background:var(--border-light);
    color:var(--ink-light);
  }
  .toolbtn.plain:hover { background:var(--border); }

  #content-area {
    flex:1;
    overflow-y:auto;
    padding:24px;
  }
  #status-bar {
    padding:6px 20px;
    background:var(--card);
    border-top:1px solid var(--border);
    font-size:0.72em;
    color:var(--ink-lighter);
    display:flex;
    justify-content:space-between;
    flex-shrink:0;
  }

  /* 内容区样式 */
  .empty-state {
    text-align:center;
    padding:80px 20px;
    color:var(--ink-lighter);
  }
  .empty-state .big-icon { font-size:3em; margin-bottom:12px; }
  .empty-state h3 { margin-bottom:8px; color:var(--ink-light); }

  .file-grid {
    display:grid;
    grid-template-columns: repeat(auto-fill, minmax(260px, 1fr));
    gap:10px;
  }
  .file-card {
    background:var(--card);
    border:1px solid var(--border);
    border-radius:8px;
    padding:14px;
    display:flex;
    align-items:center;
    gap:10px;
    cursor:pointer;
    transition: box-shadow 0.2s, border-color 0.2s;
  }
  .file-card:hover { box-shadow: 0 2px 8px rgba(0,0,0,0.08); border-color: var(--accent-gold); }
  .file-card .fc-icon { font-size:1.6em; width:32px; text-align:center; flex-shrink:0; }
  .file-card .fc-info { overflow:hidden; flex:1; min-width:0; }
  .file-card .fc-name { font-size:0.9em; font-weight:600; overflow:hidden; text-overflow:ellipsis; white-space:nowrap; }
  .file-card .fc-meta { font-size:0.72em; color:var(--ink-lighter); }
  .file-card.dir { border-left:3px solid var(--accent-gold); }

  .preview-wrapper {
    max-width:1000px;
  }
  .preview-img {
    max-width:100%;
    max-height:70vh;
    border-radius:8px;
    box-shadow:0 4px 16px rgba(0,0,0,0.1);
  }
  .preview-md {
    line-height:1.8;
    font-size:0.92em;
  }
  .preview-md h1,.preview-md h2,.preview-md h3 { margin:16px 0 8px; color:var(--ink); }
  .preview-md h1 { font-size:1.5em; border-bottom:2px solid var(--accent-gold); padding-bottom:8px; }
  .preview-md h2 { font-size:1.25em; border-bottom:1px solid var(--border); padding-bottom:4px; }
  .preview-md p { margin:8px 0; }
  .preview-md code {
    background: var(--code-bg);
    color: var(--code-fg);
    padding: 2px 6px;
    border-radius: 4px;
    font-size: 0.88em;
  }
  .preview-md pre {
    background: var(--code-bg);
    color: var(--code-fg);
    padding: 16px;
    border-radius: 8px;
    overflow-x: auto;
    font-size: 0.82em;
    line-height: 1.5;
    margin: 12px 0;
  }
  .preview-md pre code {
    background: none;
    padding: 0;
    font-size: inherit;
  }
  .preview-md table {
    border-collapse: collapse;
    width: 100%;
    margin: 12px 0;
  }
  .preview-md th, .preview-md td {
    border: 1px solid var(--border);
    padding: 6px 10px;
    text-align: left;
    font-size: 0.88em;
  }
  .preview-md th { background: var(--hover); font-weight: 600; }
  .preview-md blockquote {
    border-left: 3px solid var(--accent-gold);
    padding: 4px 16px;
    margin: 8px 0;
    color: var(--ink-light);
    background: var(--hover);
  }
  .preview-md a { color: var(--accent-gold); }

  .preview-table {
    border-collapse:collapse;
    width:100%;
    margin:8px 0;
    font-size:0.85em;
  }
  .preview-table td,.preview-table th {
    border:1px solid var(--border);
    padding:4px 8px;
  }
  .preview-table th {
    background:var(--hover);
  }

  .preview-text {
    font-family: 'SF Mono','Fira Code',monospace;
    font-size:0.82em;
    line-height:1.6;
    white-space:pre-wrap;
    word-break:break-all;
    background:var(--card);
    padding:16px;
    border-radius:8px;
    border:1px solid var(--border);
    max-height:70vh;
    overflow-y:auto;
  }

  .preview-frame {
    width:100%;
    height:80vh;
    border:1px solid var(--border);
    border-radius:8px;
  }

  /* 响应式 */
  @media (max-width:768px) {
    body { flex-direction:column; }
    #sidebar { width:100%; max-width:100%; height:40vh; resize:none; }
    #main { height:60vh; }
    .file-grid { grid-template-columns: 1fr; }
  }
</style>
</head>
<body>
<div id="sidebar">
  <div id="sidebar-header">
    <h2>📁 文件浏览器</h2>
    <p>{{ROOT_NAME}}</p>
  </div>
  <div id="tree-container">{{TREE}}</div>
  <div style="padding:8px 12px;border-top:1px solid var(--border);font-size:0.7em;color:var(--ink-lighter);text-align:center">
    满意红 Workspace · 只读浏览
  </div>
</div>
<div id="main">
  <div id="toolbar">
    <div id="breadcrumb">{{BREADCRUMB}}</div>
    <div style="flex:1"></div>
    <button class="toolbtn plain" onclick="location.reload()" title="刷新">🔄</button>
    <button class="toolbtn plain" id="dl-btn" style="display:none" onclick="downloadCurrent()">⬇️ 下载</button>
  </div>
  <div id="content-area">{{CONTENT}}</div>
  <div id="status-bar">
    <span>{{STATUS}}</span>
    <span>🕐 {{TIME}}</span>
  </div>
</div>

<script>
// ========== 状态 ==========
var currentDir = "{{CURRENT_DIR}}";
var currentFile = "{{CURRENT_FILE}}";
var currentPreview = "{{PREVIEW_TYPE}}";

// ========== 目录树交互 ==========
function toggleDir(el, dirPath) {
  event.stopPropagation();
  var children = el.parentElement.querySelector('.tree-children');
  if (!children) {
    // 懒加载
    var parentLi = el.parentElement;
    loadChildren(parentLi, dirPath, true);
    return;
  }
  children.classList.toggle('collapsed');
  el.classList.toggle('expanded');
}

function loadChildren(parentLi, dirPath, expand) {
  var existing = parentLi.querySelector('.tree-children');
  if (existing) {
    if (expand) { existing.classList.remove('collapsed'); parentLi.querySelector('.tree-toggle').classList.add('expanded'); }
    return;
  }
  fetch('/api/tree?dir=' + encodeURIComponent(dirPath))
    .then(r => r.json())
    .then(items => {
      var ul = document.createElement('ul');
      ul.className = 'tree-children';
      ul.style.listStyle = 'none';
      items.forEach(function(item) {
        var li = document.createElement('li');
        if (item.is_dir) {
          li.innerHTML = '<div class="tree-item" onclick="navigateDir(\'' + escapePath(item.path) + '\')">' +
            '<span class="tree-toggle" onclick="toggleDir(this, \'' + escapePath(item.path) + '\')">▶</span>' +
            '<span class="icon">📁</span><span class="name">' + escHtml(item.name) + '</span>' +
            '<span class="meta">' + item.modified.slice(5) + '</span></div>';
        } else {
          var icon = getFileIcon(item.ext);
          li.innerHTML = '<div class="tree-item leaf" onclick="navigateFile(\'' + escapePath(item.path) + '\')" title="' + escHtml(item.name) + '">' +
            '<span class="tree-toggle leaf">▶</span>' +
            '<span class="icon">' + icon + '</span><span class="name">' + escHtml(item.name) + '</span>' +
            '<span class="meta">' + item.size_human + '</span></div>';
        }
        ul.appendChild(li);
      });
      parentLi.appendChild(ul);
      if (expand) {
        parentLi.querySelector('.tree-toggle').classList.add('expanded');
      }
    });
}

function escHtml(s) { return s.replace(/&/g,'&amp;').replace(/</g,'&lt;').replace(/>/g,'&gt;').replace(/"/g,'&quot;'); }
function escapePath(p) { return p.replace(/'/g, "\\'"); }
function getFileIcon(ext) {
  var map = {'.md':'📝','.html':'🌐','.htm':'🌐','.py':'🐍','.js':'📜','.json':'📋','.css':'🎨','.svg':'🖼️','.png':'🖼️','.jpg':'🖼️','.jpeg':'🖼️','.gif':'🖼️','.webp':'🖼️','.pdf':'📕','.docx':'📄','.pptx':'📊','.xlsx':'📈','.txt':'📃','.csv':'📋','.yml':'⚙️','.yaml':'⚙️','.toml':'⚙️','.sh':'💻','.xml':'📋','.mp3':'🎵','.mp4':'🎬','.mov':'🎬','.zip':'📦','.gz':'📦','.tar':'📦','.ttf':'🔤','.woff':'🔤','.woff2':'🔤'};
  return map[ext] || '📄';
}

function navigateDir(path) {
  window.location.href = '/?dir=' + encodeURIComponent(path);
}
function navigateFile(path) {
  window.location.href = '/?file=' + encodeURIComponent(path);
}
function downloadCurrent() {
  if (currentFile) {
    window.open('/api/download' + currentFile, '_blank');
  }
}

// 高亮当前项
(function() {
  var currentPath = currentDir || currentFile || '/';
  var items = document.querySelectorAll('.tree-item');
  items.forEach(function(item) {
    if (item.onclick && item.onclick.toString().indexOf("'" + currentPath.replace(/\\//g,'/') + "'") > -1) {
      item.classList.add('active');
    }
  });
})();

// 下载按钮显示/隐藏
(function() {
  var btn = document.getElementById('dl-btn');
  if (currentFile && currentPreview !== 'dir') {
    btn.style.display = '';
  }
})();
</script>
</body>
</html>'''


class FileBrowserHandler(http.server.BaseHTTPRequestHandler):
    """文件浏览器请求处理器"""

    def log_message(self, format, *args):
        # 简化日志
        pass

    def _send_json(self, data, status=200):
        self.send_response(status)
        self.send_header('Content-Type', 'application/json; charset=utf-8')
        self.send_header('Cache-Control', 'no-cache')
        self.end_headers()
        self.wfile.write(json.dumps(data, ensure_ascii=False).encode('utf-8'))

    def _send_html(self, html, status=200):
        self.send_response(status)
        self.send_header('Content-Type', 'text/html; charset=utf-8')
        self.send_header('Cache-Control', 'no-cache')
        self.end_headers()
        self.wfile.write(html.encode('utf-8'))

    def _send_file_raw(self, path, content_type=None):
        try:
            with open(path, 'rb') as f:
                data = f.read()
        except OSError:
            self.send_error(404, 'File not found')
            return
        if content_type is None:
            content_type, _ = mimetypes.guess_type(str(path))
            if content_type is None:
                content_type = 'application/octet-stream'
        self.send_response(200)
        self.send_header('Content-Type', content_type)
        self.send_header('Content-Length', len(data))
        self.send_header('Cache-Control', 'max-age=3600')
        self.end_headers()
        self.wfile.write(data)

    def _send_download(self, path):
        """强制下载"""
        try:
            with open(path, 'rb') as f:
                data = f.read()
        except OSError:
            self.send_error(404)
            return
        filename = path.name
        encoded_filename = urllib.parse.quote(filename)
        self.send_response(200)
        self.send_header('Content-Type', 'application/octet-stream')
        self.send_header('Content-Disposition', f'attachment; filename="{encoded_filename}"; filename*=UTF-8\'\'{encoded_filename}')
        self.send_header('Content-Length', len(data))
        self.end_headers()
        self.wfile.write(data)

    def do_GET(self):
        parsed = urllib.parse.urlparse(self.path)
        path = parsed.path
        qs = urllib.parse.parse_qs(parsed.query)

        # API: 目录树
        if path == '/api/tree':
            dir_path = qs.get('dir', [''])[0]
            if not is_safe_path(dir_path):
                self._send_json({'error': 'Invalid path'}, 403)
                return
            tree = build_tree(dir_path)
            self._send_json(tree)
            return

        # API: 下载
        if path.startswith('/api/download'):
            file_path = path[len('/api/download'):]
            if not is_safe_path(file_path):
                self.send_error(403)
                return
            abs_path = ROOT / file_path.lstrip('/')
            if not abs_path.is_file():
                self.send_error(404)
                return
            self._send_download(abs_path)
            return

        # 主页 / 文件/目录浏览
        if path == '/' or path == '':
            dir_param = qs.get('dir', [None])[0]
            file_param = qs.get('file', [None])[0]
            self._serve_browser(dir_param, file_param)
            return

        # 静态文件回退（直接服务 workspace 内的文件）
        file_path = path.lstrip('/')
        if is_safe_path(file_path):
            abs_path = ROOT / file_path
            if abs_path.is_file():
                self._send_file_raw(abs_path)
                return

        self.send_error(404, 'Not Found')

    def _serve_browser(self, dir_path, file_path):
        """渲染浏览器主页面"""
        content_html = ''
        preview_type = 'dir'
        current_dir = ''
        current_file = ''

        if file_path:
            current_file = file_path
            if not is_safe_path(file_path):
                self.send_error(403)
                return
            abs_path = ROOT / file_path.lstrip('/')
            if abs_path.is_file():
                preview_type = 'file'
                content_html = self._preview_file(abs_path, file_path)
            else:
                self.send_error(404)
                return
        elif dir_path:
            current_dir = dir_path
            if not is_safe_path(dir_path):
                self.send_error(403)
                return
            abs_path = ROOT / dir_path.lstrip('/')
            if abs_path.is_dir():
                preview_type = 'dir'
                content_html = self._render_dir_list(abs_path, dir_path)
            else:
                self.send_error(404)
                return
        else:
            # 根目录
            current_dir = '/'
            content_html = self._render_dir_list(ROOT, '/')

        # 面包屑
        path_parts = [p for p in (current_dir or current_file).strip('/').split('/') if p]
        breadcrumb = build_breadcrumb(path_parts)

        # 目录树（根节点 + 首层展开）
        root_items = build_tree('/')
        tree_html = '<ul style="list-style:none">'
        for item in root_items:
            path_escaped = item['path'].replace("'", "\\'")
            if item['is_dir']:
                tree_html += (
                    '<li>'
                    '<div class="tree-item" onclick="navigateDir(\'{}\')">'
                    '<span class="tree-toggle" onclick="toggleDir(this, \'{}\')">▶</span>'
                    '<span class="icon">📁</span>'
                    '<span class="name">{}</span>'
                    '<span class="meta">{}</span>'
                    '</div></li>'
                ).format(path_escaped, path_escaped, item['name'], item['modified'][5:])
            else:
                icon = get_file_icon(item['ext'])
                tree_html += (
                    '<li>'
                    '<div class="tree-item leaf" onclick="navigateFile(\'{}\')">'
                    '<span class="tree-toggle leaf">▶</span>'
                    '<span class="icon">{}</span>'
                    '<span class="name">{}</span>'
                    '<span class="meta">{}</span>'
                    '</div></li>'
                ).format(path_escaped, icon, item['name'], item['size_human'])
        tree_html += '</ul>'

        # 构建页面
        html = HTML_TEMPLATE.replace('{{ROOT_NAME}}', ROOT.name)
        html = html.replace('{{TREE}}', tree_html)
        html = html.replace('{{BREADCRUMB}}', breadcrumb)
        html = html.replace('{{CONTENT}}', content_html)
        html = html.replace('{{CURRENT_DIR}}', current_dir)
        html = html.replace('{{CURRENT_FILE}}', current_file)
        html = html.replace('{{PREVIEW_TYPE}}', preview_type)
        html = html.replace('{{STATUS}}', _get_status_text(abs_path if (file_path or dir_path) else ROOT))
        html = html.replace('{{TIME}}', time.strftime('%Y-%m-%d %H:%M:%S'))

        self._send_html(html)

    def _preview_file(self, abs_path, rel_path):
        """文件预览渲染"""
        ext = abs_path.suffix.lower()
        size = abs_path.stat().st_size

        parts = [f'<div class="preview-wrapper">']
        parts.append(f'<h2 style="margin-bottom:4px;color:var(--ink)">{get_file_icon(ext)} {abs_path.name}</h2>')
        parts.append(f'<p style="color:var(--ink-lighter);font-size:0.82em;margin-bottom:16px">{_human_size(size)} · {abs_path.suffix.upper().lstrip(".")} · 路径 /{rel_path}</p>')

        try:
            if ext in IMAGE_EXTS:
                # 图片直接显示
                parts.append(f'<img src="/{rel_path}" class="preview-img" alt="{abs_path.name}">')

            elif ext == '.pdf':
                # PDF：iframe 嵌入 + 文本提取
                parts.append(f'<iframe src="/{rel_path}" class="preview-frame" title="PDF预览"></iframe>')
                if HAS_PDF:
                    parts.append('<details style="margin-top:16px"><summary style="cursor:pointer;font-weight:600;color:var(--accent-gold)">📝 文本提取</summary>')
                    parts.append(render_pdf_preview_html(abs_path))
                    parts.append('</details>')
                else:
                    parts.append('<p style="margin-top:8px;color:var(--ink-light);font-size:0.82em">PDF preview via iframe above</p>')

            elif ext in {'.docx'} and HAS_DOCX:
                parts.append(render_docx_to_html(abs_path))

            elif ext in {'.pptx'} and HAS_PPTX:
                parts.append(render_pptx_to_html(abs_path))

            elif ext in {'.xlsx'} and HAS_XLSX:
                parts.append(render_xlsx_to_html(abs_path))

            elif ext in {'.md', '.markdown'}:
                with open(abs_path, 'r', encoding='utf-8', errors='replace') as f:
                    md_text = f.read()
                parts.append('<div class="preview-md">')
                parts.append(render_md_to_html(md_text))
                parts.append('</div>')

            elif ext in {'.html', '.htm'}:
                path_quote = '/' + rel_path
                parts.append('<iframe src="' + path_quote + '" class="preview-frame" title="HTML"></iframe>')

            elif ext in {'.txt', '.py', '.js', '.css', '.json', '.yml', '.yaml', '.toml', '.xml', '.csv', '.sh', '.log', '.env', '.gitignore'}:
                with open(abs_path, 'r', encoding='utf-8', errors='replace') as f:
                    text = f.read()
                parts.append('<pre class="preview-text">' + text + '</pre>')

            else:
                icon_str = get_file_icon(ext)
                parts.append('<div class="empty-state">')
                parts.append('<div class="big-icon">' + icon_str + '</div>')
                parts.append('<h3>' + abs_path.name + '</h3>')
                parts.append('<p style="margin-bottom:16px">' + _human_size(size) + ' - preview not available</p>')
                parts.append('<a href="/api/download/' + rel_path + '" class="toolbtn" style="display:inline-block;text-decoration:none;padding:8px 20px">Download</a>')
                parts.append('</div>')
        except Exception as e:
            parts.append('<div class="empty-state"><div class="big-icon">!!</div><h3>Preview error</h3><p>' + str(e) + '</p></div>')

        parts.append('</div>')
        return '\n'.join(parts)

    def _render_dir_list(self, abs_path, rel_path):
        items = build_tree(rel_path)
        if not items:
            return '<div class="empty-state"><div class="big-icon">[FOLDER]</div><h3>Empty directory</h3></div>'
        parts = []
        parts.append('<h2 style="margin-bottom:16px;color:var(--ink)">' + abs_path.name + '</h2>')
        parts.append('<div class="file-grid">')
        import urllib.parse as _up
        for item in items:
            if item['is_dir']:
                parts.append(
                    '<a href="/?dir=' + _up.quote(item['path']) + '" style="text-decoration:none">'
                    '<div class="file-card dir">'
                    '<div class="fc-icon">&#128193;</div>'
                    '<div class="fc-info">'
                    '<div class="fc-name">' + item['name'] + '</div>'
                    '<div class="fc-meta">' + item['modified'] + '</div>'
                    '</div></div></a>'
                )
            else:
                icon = get_file_icon(item['ext'])
                parts.append(
                    '<a href="/?file=' + _up.quote(item['path']) + '" style="text-decoration:none">'
                    '<div class="file-card">'
                    '<div class="fc-icon">' + icon + '</div>'
                    '<div class="fc-info">'
                    '<div class="fc-name">' + item['name'] + '</div>'
                    '<div class="fc-meta">' + item['size_human'] + ' | ' + item['modified'] + '</div>'
                    '</div></div></a>'
                )
        parts.append('</div>')
        return '\n'.join(parts)


def get_file_icon(ext):
    icon_map = {
        '.md': '&#128221;', '.html': '&#127760;', '.htm': '&#127760;',
        '.py': '&#128013;', '.js': '&#128220;', '.json': '&#128203;',
        '.css': '&#127912;', '.svg': '&#127912;', '.png': '&#127912;',
        '.jpg': '&#127912;', '.jpeg': '&#127912;', '.gif': '&#127912;',
        '.webp': '&#127912;', '.pdf': '&#128214;', '.docx': '&#128196;',
        '.pptx': '&#128202;', '.xlsx': '&#128200;', '.txt': '&#128195;',
        '.csv': '&#128203;', '.yml': '&#9881;', '.yaml': '&#9881;',
        '.toml': '&#9881;', '.sh': '&#128187;', '.xml': '&#128203;',
        '.mp3': '&#127925;', '.mp4': '&#127916;', '.mov': '&#127916;',
        '.zip': '&#128230;', '.gz': '&#128230;', '.tar': '&#128230;'
    }
    return icon_map.get(ext, '&#128196;')


def _get_status_text(target_path):
    stat = target_path.stat()
    mtime = time.strftime('%Y-%m-%d %H:%M:%S', time.localtime(stat.st_mtime))
    if target_path.is_dir():
        return 'Directory: ' + target_path.name + ' | ' + mtime
    return 'File: ' + target_path.name + ' | ' + _human_size(stat.st_size) + ' | ' + mtime


def main():
    import socketserver
    handler = FileBrowserHandler
    print('\n' + '=' * 50)
    print('  File Browser: ' + str(ROOT.name))
    print('  URL: http://localhost:' + str(PORT))
    print('  Preview: md / html / txt / img / pdf / docx / pptx / xlsx')
    print('  Ctrl+C to stop')
    print('=' * 50 + '\n')
    with socketserver.TCPServer(('', PORT), handler) as httpd:
        httpd.allow_reuse_address = True
        try:
            httpd.serve_forever()
        except KeyboardInterrupt:
            print('\n  File browser stopped\n')


if __name__ == '__main__':
    main()
